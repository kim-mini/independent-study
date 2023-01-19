from pptx import Presentation
import copy
from pptx.util import Cm, Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import natsort

prsPath = "C:/Users/82106/Desktop/2023 01 10 -Report- Chilisin MHCD Side Crack Defect Analysis.pptx"
prs = Presentation(prsPath)

# 이미지 경로
path = 'C:/Users/82106/Desktop/121/rgb'
imgList = natsort.natsorted(os.listdir(path))
for img in imgList:
    imgCh = img.split(".")[0][-1]
    if imgCh == 'R':
        tmpStr = img
        idx = imgList.index(img)
        imgList[idx] = imgList[idx-2]
        imgList[idx-2] = tmpStr

########################################################
# 기본셋팅
# 이미지 크기 지정
scrwhidth = 4.2    # 가로
scrheight = 3.28  # 세로

# 표 정보 입력
tableleft = 2.8
tabletop = 3.42

tablewhidth = 21.5 # 가로
tableheight = 14.04 # 세로

# 가로방향으로 채워나간다
# 가로 사진개수
col = 5;
# 세로 사진개수
row = 4

########################################################

# 사진의 위치 - 행
leftlist = tableleft + ((tablewhidth / col) - scrwhidth)/2
# 사진의 위치 - 열
toplist = tabletop + ((tableheight / row) - scrheight)/2
# 이미지 높이(cm) 0 : 4면 / 1 : T1, T2
width = 0
# 이미지 넓이(cm) 0 : 4면 / 1 : T1, T2
height = 0

# 사진 간격
width1 = tablewhidth / col
height1 = tableheight / row

src = prs.slides[-1]

blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
slide = prs.slides.add_slide(blank_slide_layout)

# 첫번째 슬라이드 shape 전부 복사 후 빈 슬라이스에 붙여넣기
for shape in src.shapes:
    el = shape.element
    newel = copy.deepcopy(el)
    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

#fileList = os.listdir(path)
imgCnt = 0

# 가로 개수가 채워지면 다음 행으로 넘어간다
i = 0
j = 0
for img in imgList:
    # 열의 위치
    if i % row == 0 and imgCnt > 1:
        # 위에서 아래로
        j += 1
        leftlist += width1
        toplist = tabletop + ((tableheight / row) - scrheight)/2
    # 슬라이드 한장에 col*row개가 들어간다면 새로운 빈 슬라이드 생성
    if imgCnt > 1 and imgCnt % (col*row) == 0:
        blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
        slide = prs.slides.add_slide(blank_slide_layout)

        for shape in src.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        # 사진의 위치 - 행
        leftlist = tableleft + ((tablewhidth / col) - scrwhidth) / 2
        # 사진의 위치 - 열
        toplist = tabletop + ((tableheight / row) - scrheight) / 2

    imgPath = os.path.join(path, img)

    top = Cm(toplist)
    left = Cm(leftlist)

    toplist += height1

    width = Cm(scrwhidth)
    height = Cm(scrheight)
    # width, hegith가 없을 경우 원본 사이즈로
    # 슬라이드에 이미지 배치
    pic = slide.shapes.add_picture(imgPath, left, top, width=width, height=height)


    # 폴더안에 있는 이미지를 전부 넣었다면 중지
    if imgCnt == len(imgList)-1:
        print(imgCnt, len(imgList)-1)
        break
    imgCnt += 1
    i += 1


prs.save(prsPath)
