from pptx import Presentation
import copy
from pptx.util import Cm, Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

# ppt 경로
prs = Presentation("D:/WORK/1007 - fired chip/2021_08_ 27-Report_AVX_Fired Chip NG Sample Analysis.pptx")

# 이미지 폴더 경로
path = 'D:/WORK/1007 - fired chip/BC/13'
# 폴더들이 모여있는 경로
folderList = os.listdir(path)
imgList = os.listdir(path)

# 이미지 확장자
ext = '.tif'

# 사진의 위치 - 행
leftlist = [2.04, 5.54, 9.04, 12.54, 16.04, 19.54]
# 사진의 위치 - 열
toplist = [4.17, 6.73, 9.18, 11.73, 14.21]

imgCnt = 0
fileCnt = list()


src = prs.slides[0]
blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃 생성
slide = prs.slides.add_slide(blank_slide_layout)

# shpe name : index 형태로 dic 생성
shapeIdx = {}
for idx, value in enumerate(src.shapes):
    shapeIdx[value.name] = idx

# 첫번째 슬라이드 shape 전부 복사 후 빈 슬라이스에 붙여넣기
for shape in src.shapes:
    el = shape.element
    newel = copy.deepcopy(el)
    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

# 가로방향으로 채워나간다
# 가로로 6장이 채워지면 다음 행으로 넘어간다
i = 0
j = 0


font = '맑은 고딕 (제목)'
bnt = True

for img in imgList:
    # 열의 위치
    if img.endswith(ext):
        if j > 5 and j % 6 == 0:
            i += 1

        # 슬라이드 한장에 이미지 30개가 들어간다면 새로운 빈 슬라이드 생성
        if imgCnt > 10 and imgCnt % 30 == 0:
            blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
            slide = prs.slides.add_slide(blank_slide_layout)

        imgPath = os.path.join(path, img)

        top = Cm(toplist[i%5])
        left = Cm(leftlist[j%6])

        width = Cm(3.43)
        height = Cm(1.72)

        # width, hegith가 없을 경우 원본 사이즈로
        # 슬라이드에 이미지 배치
        pic = slide.shapes.add_picture(imgPath, left, top, width=width, height=height)


        # 폴더안에 있는 이미지를 전부 넣었다면 중지
        if imgCnt == len(imgList)-1:
            print(imgCnt, len(imgList)-1)
            break
        imgCnt += 1
        j += 1



prs.save("D:/WORK/1007 - fired chip/2021_08_ 27-Report_AVX_Fired Chip NG Sample Analysis.pptx")
