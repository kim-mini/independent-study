# 여러폴더들 돌면서 이미지 저장

from pptx import Presentation
import copy
from pptx.util import Cm, Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ppt 경로
prs = Presentation("D:/WORK/1007 - fired chip/2021_08_ 27-Report_AVX_Fired Chip NG Sample Analysis.pptx")

# 이미지 폴더 경로
path = 'D:/WORK/1007 - fired chip/BC'
# 폴더들이 모여있는 경로
folderList = os.listdir(path)

# 이미지 확장자
ext = '.png'

# 사진의 위치 - 행
leftlist = [2.04, 5.54, 9.04, 12.54, 16.04, 19.54]
# 사진의 위치 - 열
toplist = [4.17, 6.73, 9.18, 11.73, 14.21]

imgCnt = 0
fileCnt = list()


src = prs.slides[0]
blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃 생성
slide = prs.slides.add_slide(blank_slide_layout)

# 첫번째 슬라이드 shape 전부 복사 후 빈 슬라이스에 붙여넣기
for shape in src.shapes:
    el = shape.element
    newel = copy.deepcopy(el)
    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

# shpe name : index 형태로 dic 생성
shapeIdx = {}
for idx, value in enumerate(slide.shapes):
    shapeIdx[value.name] = idx

# 가로방향으로 채워나간다
# 가로로 6장이 채워지면 다음 행으로 넘어간다
i = 0
j = 0


font = '맑은 고딕 (제목)'
bnt = True

number = 1

for folder in folderList:
    folderPath = os.path.join(path, folder)
    imgList = os.listdir(folderPath)

    for img in imgList:
        # 설정한 확장자를 가진 파일들만 작업을 시작함
        if img.endswith(ext):
            if bnt:
                slide = prs.slides.add_slide(blank_slide_layout)

                for shape in src.shapes:
                    el = shape.element
                    newel = copy.deepcopy(el)
                    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

                folderText = folder.split('.')[1]
                # name id shape 선택
                shapeName = slide.shapes[shapeIdx['name']]
                tf = shapeName.text_frame
                tf.clear()

                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.RIGHT
                run = para.add_run()
                run.text = f'{number}) {folderText} #1 ~ #'
                font = run.font
                font.size = Pt(14)
                font.name = '맑은 고딕 (제목)'
                font.bold = True
                font.color.rgb = RGBColor(0, 0, 0)
                number += 1
                bnt = False

            # 열의 위치
            if j > 5 and j % 6 == 0:
                i += 1

            # 슬라이드 한장에 이미지 30개가 들어간다면 새로운 빈 슬라이드 생성
            if imgCnt > 10 and imgCnt % 30 == 0:
                break


            imgPath = os.path.join(folderPath, img)

            top = Cm(toplist[i%5])
            left = Cm(leftlist[j%6])
            # 4면 가로세로 길이 정의

            width = Cm(3.43)
            height = Cm(1.72)

            # width, height가 없을 경우 원본 사이즈로
            # 슬라이드에 이미지 배치
            pic = slide.shapes.add_picture(imgPath, left, top, width=width, height=height)


            # 폴더안에 있는 이미지를 전부 넣었다면 중지
            if imgCnt == len(imgList)-1:
                print('hey : ', imgCnt, len(imgList) - 1)
                break

            imgCnt += 1
            j += 1

        else:
            fileCnt.append(0)
            if imgCnt == len(imgList) - (1 + fileCnt.count(0)):
                break


    bnt = True
    imgCnt, i, j = 0, 0, 0

prs.save("D:/WORK/1007 - fired chip/2021_08_ 27-Report_AVX_Fired Chip NG Sample Analysis.pptx")
