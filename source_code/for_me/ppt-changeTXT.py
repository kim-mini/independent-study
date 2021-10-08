# ppt 내부 텍스트 바꾸기

from pptx import Presentation
from pptx.util import Pt
import re

prs = Presentation("C:/Users/82106/Desktop/1/2. 20210607-AVX-WVS-MCI6CH Manual System set-up & Maintenance.pptx")
result = []
cnt = 1

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            if '[Fig' in paragraph.text:
                # 폰트이름
                paragraph.font.name = 'Noto Sans'
                # 폰트 사이즈
                paragraph.font.size = Pt(8)
                # [Fig#15 abcdefg0]
                p = paragraph.text.split(' ', maxsplit=1)
                endP = p[1] # 뒷부분
                startP = p[0].split('#')[0] # 숫자빼고 앞부분만
                text = startP + '#' + str(cnt)
                text = text + ' ' + endP
                paragraph.text = text
                result.append(paragraph.text)
                cnt += 1


print(result)
prs.save("C:/Users/82106/Desktop/1/2. 20210607-AVX-WVS-MCI6CH Manual System set-up & Maintenance.pptx")
